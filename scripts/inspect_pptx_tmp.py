#!/usr/bin/env python3
from pptx import Presentation
from pathlib import Path

ppt = Path('/rds/general/user/moa324/home/projects/cyberwheel/research_docs/cyberwheel_slides.pptx')
if not ppt.exists():
    print('PPTX not found:', ppt)
    raise SystemExit(1)

prs = Presentation(str(ppt))
print('Slides:', len(prs.slides))
for i, slide in enumerate(prs.slides, start=1):
    title = slide.shapes.title.text if slide.shapes.title is not None else '<no title>'
    imgs = [shape for shape in slide.shapes if shape.shape_type == 13]  # 13 == PICTURE
    print(f'Slide {i}: title={title!r}, images={len(imgs)}')
    if imgs:
        for j,im in enumerate(imgs, start=1):
            print('  image', j, 'size', im.width, 'x', im.height)
