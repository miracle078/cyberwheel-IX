#!/usr/bin/env python3
"""Extract 'Problem Formulation and Environment' from LaTeX and append slides.

Workflow:
 - Extract section from `research_docs/original.tex` (or original copy)
 - Split into headings/paragraphs and math fragments
 - Render math fragments to SVG and PNG (using matplotlib mathtext)
 - Append slides to target PPTX (`research_docs/cyberwheel_slides_vector_media.pptx`) embedding the PNGs
 - Save PPTX as `research_docs/cyberwheel_slides_with_formalisation.pptx`
 - Run the vector injector to replace PNGs with SVG parts when possible

Notes: Requires python-pptx and matplotlib. If matplotlib isn't available the script will still append plain text slides but won't render math images.
"""

import re
import os
import sys
import zipfile
import subprocess

from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
FIGDIR = ROOT / 'research_docs' / 'figures'
TEX_CANDS = [ROOT / 'research_docs' / 'original.tex', ROOT / 'research_docs' / 'original copy.tex']
TARGET_PPTX = ROOT / 'research_docs' / 'cyberwheel_slides_vector_media.pptx'
OUT_PPTX = ROOT / 'research_docs' / 'cyberwheel_slides_with_formalisation.pptx'


def find_tex():
    for p in TEX_CANDS:
        if p.exists():
            return p
    raise FileNotFoundError('No TeX source found in expected locations')


def extract_section(tex_path, section_title='Problem Formulation and Environment'):
    text = tex_path.read_text(encoding='utf-8')
    # naive extraction: find \section{...} and next \section
    pattern = re.compile(r"\\section\{\s*" + re.escape(section_title) + r"\s*\}(.*?)(?=\\section\{|\\end\{document\}|$)", re.S)
    m = pattern.search(text)
    if not m:
        # fallback: try without exact match
        pattern2 = re.compile(r"\\section\{Problem Formulation.*?\}(.*?)(?=\\section\{|$)", re.S)
        m = pattern2.search(text)
    if not m:
        raise SystemExit('Could not find Problem Formulation section in TeX')
    return m.group(1).strip()


def chunk_to_slides(section_text):
    # split by subsections or blank-lines into slide-sized chunks
    parts = re.split(r"\\subsection\{|\\subsubsection\{|\\paragraph\{|\\subsection\*\{|\\subsubsection\*\{", section_text)
    # simplistic: also split on double newlines
    if len(parts) <= 1:
        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", section_text) if p.strip()]
    else:
        paragraphs = [p.strip() for p in parts if p.strip()]
    slides = []
    for p in paragraphs:
        # remove remaining TeX commands mildly
        clean = re.sub(r"\\cwfigure\{.*?\}\{.*?\}\{.*?\}", '', p, flags=re.S)
        slides.append(clean)
    return slides


def detect_math_fragments(text):
    # find display math and inline math
    fragments = []
    # display $$...$$ and \[...\]
    for m in re.finditer(r"(\\\[.*?\\\]|\\begin\{equation\}.*?\\end\{equation\}|\$\$.*?\$\$)", text, re.S):
        fragments.append((m.start(), m.end(), m.group(0)))
    # inline $...$
    for m in re.finditer(r"(?<!\$)\$(?!\$)(.+?)(?<!\$)\$(?!\$)", text, re.S):
        fragments.append((m.start(), m.end(), m.group(0)))
    fragments.sort()
    return fragments


def render_math_to_svg_png(math_tex, out_svg, out_png):
    # Try using matplotlib.mathtext to render to SVG and PNG
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        from matplotlib import rc
        fig = plt.figure()
        fig.text(0, 0, f"${math_tex.strip().strip('$').strip()}$", fontsize=14)
        # save SVG and PNG
        fig.savefig(str(out_svg), format='svg', bbox_inches='tight', pad_inches=0.02)
        fig.savefig(str(out_png), format='png', dpi=300, bbox_inches='tight', pad_inches=0.02)
        plt.close(fig)
        return True
    except Exception as e:
        print('matplotlib rendering failed:', e)
        return False


def next_media_index(pptx_path):
    with zipfile.ZipFile(pptx_path, 'r') as z:
        media = [n for n in z.namelist() if n.startswith('ppt/media/image')]
        nums = [int(re.search(r'image(\d+)\.', m).group(1)) for m in media]
        return max(nums) + 1 if nums else 1


def append_slides(pptx_in, pptx_out, slides_content):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as e:
        raise SystemExit('python-pptx is required: ' + str(e))

    prs = Presentation(str(pptx_in))
    start_idx = next_media_index(pptx_in)
    os.makedirs(FIGDIR, exist_ok=True)
    media_idx = start_idx

    for i, content in enumerate(slides_content):
        # create title+content slide
        slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = 'Formalisation' if i == 0 else f'Formalisation (cont.)'

        body = slide.shapes.placeholders[1] if len(slide.shapes.placeholders) > 1 else None
        text_frame = body.text_frame if body is not None else None

        # find math fragments and replace them with images sequentially
        frags = detect_math_fragments(content)
        cursor = 0
        if not frags:
            # plain paragraph
            if text_frame:
                text_frame.text = re.sub(r'\\[a-zA-Z]+\*?\{.*?\}', '', content)
        else:
            # build text with images appended below
            plain = re.sub(r'\$\$.*?\$\$|\\\[.*?\\\]|\$.*?\$|\\begin\{equation\}.*?\\end\{equation\}', ' [MATH] ', content, flags=re.S)
            if text_frame:
                text_frame.text = re.sub(r'\\[a-zA-Z]+\*?\{.*?\}', '', plain)

            for s,e,frag in frags:
                math = frag
                # clean delimiters
                math_clean = re.sub(r'^\\\[|\\\]$|^\$\$|\$\$$|^\$|\$$|\\begin\{equation\}|\\end\{equation\}', '', math, flags=re.S).strip()
                svg_path = FIGDIR / f'fig_{media_idx}.svg'
                png_path = FIGDIR / f'fig_{media_idx}.png'
                ok = render_math_to_svg_png(math_clean, svg_path, png_path)
                if not ok:
                    # fallback: write simple text
                    if text_frame:
                        p = text_frame.add_paragraph()
                        p.text = math_clean
                else:
                    # insert picture (PNG) into slide
                    left = Inches(1)
                    top = Inches(2)
                    height = Inches(1.5)
                    try:
                        slide.shapes.add_picture(str(png_path), left, top, height=height)
                    except Exception:
                        # if can't add picture, append text path
                        if text_frame:
                            p = text_frame.add_paragraph()
                            p.text = f'[Equation image: {png_path.name}]'
                    media_idx += 1

    prs.save(str(pptx_out))
    return pptx_out


def run_injector(pptx_path):
    injector = ROOT / 'research_docs' / 'scripts' / 'inject_vectors_into_pptx.py'
    out = str(pptx_path).replace('.pptx', '_vectors_injected.pptx')
    cmd = [sys.executable, str(injector), '--input', str(pptx_path), '--out', out, '--figdir', str(FIGDIR)]
    print('Running injector:', ' '.join(cmd))
    subprocess.check_call(cmd)
    return out


def main():
    tex = find_tex()
    section = extract_section(tex)
    slides = chunk_to_slides(section)
    if not slides:
        raise SystemExit('No content parsed into slides')
    print(f'Parsed {len(slides)} slide chunks from TeX')
    appended = append_slides(TARGET_PPTX, OUT_PPTX, slides)
    print('Appended slides to', appended)
    injected = run_injector(OUT_PPTX)
    print('Vector-injected PPTX written to', injected)
    print('\nChecklist:\n - Extracted section from %s\n - Rendered math to SVG+PNG where possible into %s\n - Appended slides to %s\n - Replaced PNGs with SVGs/EMFs in %s' % (tex, FIGDIR, OUT_PPTX, injected))


if __name__ == '__main__':
    main()
