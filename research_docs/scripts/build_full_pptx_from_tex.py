"""
Recreate the full Beamer slides as native PowerPoint slides.
This is a best-effort parser: it extracts `\begin{frame}` blocks, `\frametitle{}`, simple `\framesubtitle{}`, `itemize` lists, `enumerate`, and plain paragraphs.
It embeds vector images if `research_docs/figures/fig_N.emf` (or .wmf/.svg/.png) exist for the slide.

Usage: python3 build_full_pptx_from_tex.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import re

REPO_ROOT = Path(__file__).resolve().parents[2]
TEX = REPO_ROOT / 'research_docs' / 'cyberwheel_slides_final.tex'
FIG_DIR = REPO_ROOT / 'research_docs' / 'figures'
OUT = REPO_ROOT / 'research_docs' / 'cyberwheel_slides_recreated_emf.pptx'


def parse_frames(tex_path: Path):
    text = tex_path.read_text()
    frames = []
    # split on \begin{frame} but keep content
    parts = re.split(r'\\begin\{frame\}', text)
    for p in parts[1:]:
        content = p
        title = ''
        subtitle = ''
        if '\\frametitle{' in content:
            title = re.search(r'\\frametitle\{([^}]*)\}', content)
            title = title.group(1).strip() if title else ''
        if '\\framesubtitle{' in content:
            subtitle = re.search(r'\\framesubtitle\{([^}]*)\}', content)
            subtitle = subtitle.group(1).strip() if subtitle else ''
        # itemize/enumerate
        items = []
        if '\\begin{itemize}' in content:
            b = re.search(r'\\begin\{itemize\}(.+?)\\end\{itemize\}', content, re.S)
            if b:
                for m in re.finditer(r'\\item\s*([^\\]+)', b.group(1)):
                    items.append(m.group(1).strip())
        # paragraphs: naive - lines not starting with % or \n etc
        paragraphs = []
        # remove environments for simplicity
        clean = re.sub(r'\\begin\{[^}]+\}.*?\\end\{[^}]+\}', '', content, flags=re.S)
        for line in clean.splitlines():
            line=line.strip()
            if not line: continue
            if line.startswith('%'): continue
            if line.startswith('\\') or line.startswith('{') or line.startswith('}'): continue
            # skip frametitle/framesubtitle lines
            if line.startswith('\\frametitle') or line.startswith('\\framesubtitle'): continue
            paragraphs.append(line)
        frames.append({'title': title, 'subtitle': subtitle, 'items': items, 'paragraphs': paragraphs, 'raw': content})
    return frames


def find_media_for_index(i):
    # prefer emf/wmf/svg/png
    for ext in ('emf','wmf','svg','png'):
        p = FIG_DIR / f'fig_{i+1}.{ext}'
        if p.exists():
            return p
    return None


def add_bulleted_text(slide, left, top, width, height, items):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    for idx, it in enumerate(items):
        p = tf.add_paragraph() if idx>0 else tf.paragraphs[0]
        p.text = it
        p.level = 0
        p.font.size = Pt(18)


if __name__ == '__main__':
    frames = parse_frames(TEX)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, f in enumerate(frames):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # title
        if f['title']:
            slide.shapes.title.text = f['title']
        # subtitle as smaller textbox
        if f['subtitle']:
            sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
            tf = sub.text_frame
            tf.text = f['subtitle']
            tf.paragraphs[0].font.size = Pt(14)
        # paragraphs and items
        if f['items']:
            add_bulleted_text(slide, Inches(0.5), Inches(1.3), Inches(4.2), Inches(5), f['items'])
        elif f['paragraphs']:
            txt = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8.5), Inches(5))
            tf = txt.text_frame
            tf.text = '\n'.join(f['paragraphs'][:4])
            for p in tf.paragraphs:
                p.font.size = Pt(18)
        # add media if exists
        media = find_media_for_index(i)
        if media:
            # left or right placement heuristic
            slide.shapes.add_picture(str(media), Inches(5.2), Inches(1.1), width=Inches(4.8))

    prs.save(OUT)
    print('Wrote', OUT)
