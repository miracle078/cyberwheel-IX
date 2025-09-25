#!/usr/bin/env python3
r"""
Build PPTX from the original Beamer source by parsing \begin{frame}...\end{frame} blocks.
Uses the exported SVGs in research_docs/figures/ (fig_1.svg, fig_2.svg, ...) for tikzpicture blocks
and embeds PNG fallbacks converted via cairosvg for reliable insertion.

Output: research_docs/cyberwheel_slides_from_original.pptx
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

tex = Path('/rds/general/user/moa324/home/projects/cyberwheel/research_docs/cyberwheel_slides_final.tex')
if not tex.exists():
    print('Original tex not found:', tex)
    sys.exit(1)

fig_dir = tex.parent / 'figures'
# list available figure svgs in order
svgs = sorted(fig_dir.glob('fig_*.svg'))
if not svgs:
    print('No exported SVGs found in', fig_dir)

text = tex.read_text(encoding='utf-8')
# capture frames including optional title in \begin{frame}{Title} or \begin{frame}[...]{Title}
frame_re = re.compile(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{([^}]*)\})?(.*?)\\end\{frame\}", re.S)
frames = frame_re.findall(text)
print('Found', len(frames), 'frames')

prs = Presentation()
# set slide size to standard 4:3 (10in x 7.5in)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

# helpers
itemize_re = re.compile(r"\\begin\{itemize\}(.*?)\\end\{itemize\}", re.S)
item_re = re.compile(r"\\item\s*(.*?)($|\\\\|\\item)", re.S)

svg_idx = 0

# ensure cairosvg available
try:
    import cairosvg
    cairosvg_available = True
except Exception:
    cairosvg_available = False

for i, (title, body) in enumerate(frames, start=1):
    slide = prs.slides.add_slide(blank_layout)
    # Title
    frame_title = title.strip() if title and title.strip() else None
    # also try \frametitle{} inside body
    if not frame_title:
        m = re.search(r"\\frametitle\{([^}]*)\}", body)
        if m:
            frame_title = m.group(1).strip()
    if frame_title:
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.6))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = frame_title
        p.font.size = Pt(26)

    # Extract bullets (first itemize block)
    bullets = []
    m = itemize_re.search(body)
    if m:
        block = m.group(1)
        for li in item_re.findall(block):
            bullets.append(re.sub(r"\\[a-zA-Z]+\{([^}]*)\}", r"\1", li[0]).strip())

    if bullets:
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.95), Inches(4.6), Inches(5))
        tf = tx.text_frame
        tf.word_wrap = True
        for idx_b, b in enumerate(bullets):
            p = tf.add_paragraph() if idx_b>0 else tf.paragraphs[0]
            p.text = b
            p.level = 0
            p.font.size = Pt(14)

    # Detect tikzpicture blocks in the frame body
    tikz_re = re.compile(r"\\begin\{tikzpicture\}.*?\\end\{tikzpicture\}", re.S)
    tikz_blocks = tikz_re.findall(body)
    if tikz_blocks:
        for tb in tikz_blocks:
            if svg_idx < len(svgs):
                svg_path = svgs[svg_idx]
                svg_idx += 1
                # attempt to embed SVG directly; python-pptx may not accept SVG, so convert to PNG
                png_path = svg_path.with_suffix('.png')
                try:
                    # convert SVG->PNG for embedding
                    if not png_path.exists():
                        if not cairosvg_available:
                            print('cairosvg not available, installing...')
                            import subprocess
                            subprocess.check_call(['conda','install','-c','conda-forge','cairosvg','-y'])
                            import importlib
                            try:
                                cairosvg = importlib.import_module('cairosvg')
                                cairosvg_available = True
                            except Exception:
                                cairosvg_available = False
                        if cairosvg_available:
                            cairosvg.svg2png(url=str(svg_path), write_to=str(png_path), dpi=300)
                            print('Converted', svg_path.name, '->', png_path.name)
                        else:
                            print('Cannot convert SVG; skip image embedding for', svg_path.name)
                            continue
                    # placement
                    if bullets:
                        left = Inches(5.2); top = Inches(0.95); width = Inches(4.0)
                    else:
                        left = Inches(0.4); top = Inches(0.95); width = Inches(9.0)
                    slide.shapes.add_picture(str(png_path), left, top, width=width)
                    print('Embedded image for frame', i, 'from', svg_path.name)
                except Exception as e:
                    print('Failed embedding image', svg_path, '->', e)
            else:
                print('No remaining SVGs to map for tikz block in frame', i)

# Save output
out = tex.parent / 'cyberwheel_slides_from_original_4x3.pptx'
prs.save(str(out))
print('Wrote', out, 'with', len(prs.slides), 'slides (4:3)')
