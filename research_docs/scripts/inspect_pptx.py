#!/usr/bin/env python3
"""
Basic inspector for a PPTX to list slides, titles, and image shapes.
Usage: python3 inspect_pptx.py <pptx-file>
"""
import sys
from pptx import Presentation

if len(sys.argv) < 2:
    print('Usage: inspect_pptx.py <pptx-file>')
    sys.exit(1)

prs = Presentation(sys.argv[1])
print('Slides:', len(prs.slides))
for i,slide in enumerate(prs.slides, start=1):
    title = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                title = txt
                break
    images = [s for s in slide.shapes if s.shape_type==13]  # 13 is picture
    print(f'Slide {i}: title="{title}" images={len(images)}')
