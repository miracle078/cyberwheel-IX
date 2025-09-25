#!/usr/bin/env python3
r"""
Build a PPTX programmatically from a LaTeX file where TikZ figures were replaced by
\includegraphics{figures/fig_<n>.svg}. Embeds SVGs directly if python-pptx allows; otherwise
falls back to converting SVG->PNG via cairosvg.

Writes: research_docs/cyberwheel_slides_programmatic.pptx
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

tex_path = Path('/rds/general/user/moa324/home/projects/cyberwheel/research_docs/cyberwheel_slides_for_pptx.tex')
out_pptx = tex_path.parent / 'cyberwheel_slides_programmatic.pptx'
if not tex_path.exists():
    print('Input tex not found:', tex_path)
    sys.exit(1)

text = tex_path.read_text(encoding='utf-8')
# Regex to capture optional options, optional title, and body
frame_re = re.compile(r"\\begin\{frame\}(?:\[(.*?)\])?(?:\{([^}]*)\})?(.*?)\\end\{frame\}", re.S)
frames = frame_re.findall(text)
print('Found', len(frames), 'frames')

prs = Presentation()
blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

# helper to strip simple LaTeX items
def strip_latex(s):
    s = re.sub(r'\\textbf\{([^}]*)\}', r'\1', s)
    s = re.sub(r'\\emph\{([^}]*)\}', r'\1', s)
    s = re.sub(r'\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{[^}]*\})?', '', s)
    s = s.replace('\\', '\n')
    return s.strip()

# try embedding SVG directly; if fails later, we'll convert
cairosvg_available = False
try:
    import cairosvg
    cairosvg_available = True
except Exception:
    cairosvg_available = False

for idx, (options, title, body) in enumerate(frames, start=1):
    slide = prs.slides.add_slide(blank_layout)
    # Title
    frame_title = title.strip() if title else ''
    if frame_title:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = frame_title
        p.font.size = Pt(28)

    # Find includegraphics
    imgs = re.findall(r"\\includegraphics\[.*?\]\{([^}]+)\}", body)
    # Find itemize bullets
    bullets = []
    for it in re.finditer(r"\\begin\{itemize\}(.*?)\\end\{itemize\}", body, re.S):
        block = it.group(1)
        for li in re.findall(r"\\item\s+(.*?)($|\\\\|\\item)", block, re.S):
            bullets.append(strip_latex(li[0]).replace('\n',' ').strip())

    # Add bullets textbox on left if present
    left = Inches(0.5)
    top = Inches(1.0)
    if bullets:
        tx = slide.shapes.add_textbox(left, top, Inches(4.5), Inches(5))
        tf = tx.text_frame
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
            p.text = b
            p.level = 0
            p.font.size = Pt(14)

    # Insert first image (if any)
    if imgs:
        img_rel = imgs[0]
        # resolve path
        img_path = (tex_path.parent / img_rel).resolve()
        if not img_path.exists():
            print('Image not found:', img_path)
        else:
            # target placement: right side, or full width if no bullets
            if bullets:
                pic_left = Inches(5.2)
                pic_top = Inches(1.0)
                pic_width = Inches(4.0)
            else:
                pic_left = Inches(0.5)
                pic_top = Inches(1.0)
                pic_width = Inches(9.0)

            try:
                slide.shapes.add_picture(str(img_path), pic_left, pic_top, width=pic_width)
                print('Embedded', img_path.name, 'into slide', idx)
            except Exception as e:
                print('Failed to embed SVG directly:', e)
                if not cairosvg_available:
                    print('Attempting to install cairosvg via conda...')
                    import subprocess
                    subprocess.check_call(['conda','install','-c','conda-forge','cairosvg','-y'])
                        import importlib
                        try:
                            cairosvg = importlib.import_module('cairosvg')
                            cairosvg_available = True
                        except Exception:
                            try:
                                # Try direct import as a last resort
                                import cairosvg
                                cairosvg_available = True
                            except Exception:
                                cairosvg_available = False
                if cairosvg_available:
                    png_out = img_path.with_suffix('.png')
                    print('Converting', img_path, '->', png_out)
                        cairosvg.svg2png(url=str(img_path), write_to=str(png_out))
                    slide.shapes.add_picture(str(png_out), pic_left, pic_top, width=pic_width)
                    print('Embedded PNG fallback for', img_path.name, 'into slide', idx)
                else:
                    print('No SVG conversion available; skipping image for slide', idx)

# Save
prs.save(str(out_pptx))
print('Wrote', out_pptx)
