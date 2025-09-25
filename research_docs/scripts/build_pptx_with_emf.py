"""
Clean EMF-based PPTX builder.
Embeds fig_N.emf when available; falls back to fig_N.png or fig_N.svg->png.
"""
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
from cairosvg import svg2png

# Repository root: .../projects/cyberwheel
REPO_ROOT = Path(__file__).resolve().parents[2]
FIG_DIR = REPO_ROOT / 'research_docs' / 'figures'
OUT = REPO_ROOT / 'research_docs' / 'cyberwheel_slides_from_original_emf.pptx'


def extract_frames(texpath: Path):
    text = texpath.read_text()
    frames = []
    parts = text.split('\\begin{frame}')
    for p in parts[1:]:
        title = ''
        if '\\frametitle{' in p:
            title = p.split('\\frametitle{', 1)[1].split('}', 1)[0].strip()
        items = []
        if '\\begin{itemize}' in p:
            body = p.split('\\begin{itemize}', 1)[1].split('\\end{itemize}', 1)[0]
            for line in body.splitlines():
                line = line.strip()
                if line.startswith('\\item'):
                    items.append(line.replace('\\item', '').strip())
        frames.append({'title': title, 'items': items, 'content': p})
    return frames


def ensure_png_for_svg(svg_path: Path, png_path: Path):
    if png_path.exists():
        return
    try:
        svg2png(url=str(svg_path), write_to=str(png_path), scale=2.0)
    except Exception:
        pass


def add_slide_with_media(prs: Presentation, title: str, items: list, media_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title or ''
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(4.0)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = it
        p.level = 0

    if media_path and media_path.exists() and media_path.is_file():
        slide.shapes.add_picture(str(media_path), Inches(5.2), Inches(1.1), width=Inches(4.8))


if __name__ == '__main__':
    TEX = REPO_ROOT / 'research_docs' / 'cyberwheel_slides_final.tex'
    if not TEX.exists():
        raise SystemExit(f'TeX file not found: {TEX}')
    frames = extract_frames(TEX)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, f in enumerate(frames):
        emf = FIG_DIR / f'fig_{i+1}.emf'
        svg = FIG_DIR / f'fig_{i+1}.svg'
        png = FIG_DIR / f'fig_{i+1}.png'
        media = None
        if emf.exists():
            media = emf
        elif png.exists():
            media = png
        elif svg.exists():
            ensure_png_for_svg(svg, png)
            if png.exists():
                media = png

        add_slide_with_media(prs, f['title'], f['items'], media if media else Path(''))

    prs.save(OUT)
    print('Wrote', OUT)
